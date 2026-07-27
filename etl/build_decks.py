"""
build_decks.py -- assembles the two Week-6 decks (.pptx) from docs/fig/ figures + results.

Styled to match the user's own deck aesthetic: minimal (plain black title top-left, no color
bars/kickers), figure-forward, dash bullets in a "Term: explanation" format with inline
citations. Re-runnable: regenerate figures, re-run this.

  1. docs/Bob_Week6_Update.pptx   -- client update for Bob/Adam (tied to their Wk5 asks)
  2. docs/MSDS_498_Midterm.pptx   -- Week-6 midterm presentation (methods/rigor)

Run:  python etl/build_deck_figures.py && python etl/build_decks.py

Every slide carries speaker notes (visible in PowerPoint's presenter view).

Two conventions, deliberate:
  * No internal version numbers on any slide. The index has been through several revisions;
    that is repo bookkeeping, not audience content, and "v4" on a slide only invites
    "what happened to the other three?". Figures and text are all cut from the same build so
    they agree -- which is the reason the versioning mattered in the first place.
  * The deliverable is described as a school-level measure, never a student-level one. There is
    no admissions outcome in this data and therefore no admit model; saying otherwise invites
    the one question the project cannot answer.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

try:
    from PIL import Image
    def _img_wh(p):
        with Image.open(p) as im:
            return im.size
except Exception:
    def _img_wh(p):
        return (1600, 900)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "docs", "fig")

# Set once, appears on the midterm title slide.
TEAM = "Capstone Team  ·  July 2026   (add team member names)"

BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x55, 0x55)
FONT = "Arial"
SW, SH = Inches(13.333), Inches(7.5)


def _fit(path, max_w, max_h):
    w, h = _img_wh(path)
    ar = w / h
    if ar >= max_w / max_h:
        return max_w, max_w / ar
    return max_h * ar, max_h


def _set(run, size, color=BLACK, bold=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.name = FONT


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def title_slide(prs, title, subtitle, foot, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(12), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = title; _set(r, 38, BLACK, True)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle; _set(r2, 19, GREY)
    fb = s.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(12), Inches(0.5))
    fr = fb.text_frame.paragraphs[0].add_run(); fr.text = foot; _set(fr, 12, GREY)
    _notes(s, notes)
    return s


def _bullets(frame, bullets):
    for i, b in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        sub = b.startswith("- ")
        if sub:
            b = b[2:]
        r = p.add_run()
        r.text = ("     ·  " if sub else "-  ") + b
        _set(r, 13 if sub else 15.5, GREY if sub else BLACK)
        p.space_after = Pt(9 if not sub else 5)
        p.line_spacing = 1.05


def content_slide(prs, title, bullets, image=None, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(1.0))
    tr = tb.text_frame.paragraphs[0].add_run(); tr.text = title; _set(tr, 30, BLACK, True)
    tb.text_frame.word_wrap = True

    path = os.path.join(FIG, image) if image else None
    if path and os.path.exists(path):
        # figure large on the left, a few key-point bullets on the right
        w, h = _fit(path, 8.0, 5.4)
        top = Inches(1.55) + Inches((5.4 - h) / 2)
        s.shapes.add_picture(path, Inches(0.45), top, width=Inches(w), height=Inches(h))
        bx = s.shapes.add_textbox(Inches(8.7), Inches(1.6), Inches(4.35), Inches(5.4))
    else:
        if image:
            print(f"  !! missing figure {image} -- run build_deck_figures.py first")
        bx = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(5.3))
    bx.text_frame.word_wrap = True
    _bullets(bx.text_frame, bullets)
    _notes(s, notes)
    return s


def new_deck():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    return prs


# ===================================================================== BOB DECK
def build_bob():
    prs = new_deck()
    title_slide(prs,
        "High School Rigor Methodology",
        "Week-6 Client Update  ·  for Bob & Adam, NU Admissions",
        "MSDS 498 Capstone  ·  July 2026",
        notes="Frame the meeting: everything here answers something they asked for in Week 5. "
              "Close on the per-course AP data ask -- that is what we need from them.")

    content_slide(prs, "Where we are", [
        "Goal: a transparent, school-level measure of academic rigor, from public data, that NU can tie to its own internal data.",
        "Built from your Week-5 asks: AP exam scores, the low-offering/high-scores idea, coverage transparency, STEM, graduation rate.",
        "Validated: the tiers track real outcomes, and are not just a proxy for socioeconomic status.",
        "Scope: we describe the school an applicant came from — not the applicant.",
    ], notes="The last bullet is the guardrail. If they ask 'can this rank students?', the "
             "answer is no by design: this gives a reader context for a transcript.")

    content_slide(prs, "How we measure rigor — an index, not a black box", [
        "Composite index: a weighted average of five standardized components → one rigor score.",
        "Decomposable: every tier traces back to the features that produced it (the Landscape lesson).",
        "Performance-weighted: AP exam scores, not just what's offered (Geiser & Santelices, 2004).",
        "No imputation: a school is scored on the components it actually has, or not scored at all.",
    ], image="index_schematic.png",
       notes="Landscape's failure mode was opacity. Ours is auditable: any tier can be walked "
             "back to its inputs, which is what lets them defend a decision.")

    content_slide(prs, "Your ask: not equal buckets", [
        "Method: natural breaks — cuts fall at the real gaps in the distribution, not at fixed percentiles.",
        "Result: 'Most Demanding' is 295 schools, not a forced top 20%.",
        "Tier sizes vary because schools genuinely cluster — 8,905 land in the middle.",
        "Alternative kept on the shelf: equal-fifths, if you ever want fixed-size tiers.",
    ], image="tier_cutpoints.png",
       notes="This is a direct answer to Bob's Week-5 pushback on equal buckets. The two schemes "
             "agree on only about half of schools, so the choice is consequential -- worth saying.")

    content_slide(prs, "What a school looks like when it comes out", [
        "Deliverable: per school — a tier, a score, a percentile, and the components behind it.",
        "Auditable: a counselor or reader can see exactly why a school landed where it did.",
        "Joinable: keyed on CEEB, so it drops straight into your existing systems.",
        "Coverage note: this school scored on 4 of 5 components — we show that, we don't hide it.",
    ], image="school_profile.png",
       notes="Show the product, not just the method. Expect them to name a school they know and "
             "ask where it lands -- New Trier is Very Demanding, 96th percentile, and that is a "
             "reasonable-sounding answer to defend.")

    content_slide(prs, "Your idea: low AP offering / high AP scores", [
        "Selective & effective: few APs offered, high exam scores — 1,597 schools.",
        "The catch: only 3 reach the top tier; a thin catalog drags the composite down.",
        "So: these 'do a lot with little' schools are hidden by the tier alone.",
        "Use: a recruiting signal — ships alongside the tier, not folded into it.",
    ], image="ap_efficiency.png",
       notes="This came directly from Bob. Worth crediting him. The honest finding is that the "
             "additive index cannot surface these schools, so we ship the flag separately.")

    content_slide(prs, "Data coverage — public vs private", [
        "Not missing-at-random: reported per sector.",
        "CRDC: AP participation, testtaker rate, grad rate — public-only by law → ~0% private.",
        "NU analytics: AP/SAT scores track your recruiting universe (~35%).",
        "Consequence: private-school tiers rest on the thin NU block + IB — a documented blind spot.",
    ], image="coverage_by_sector.png",
       notes="Answers their Week-5 question about the private-school number. Do not soften this: "
             "the blind spot is structural, and naming it is what makes the rest credible.")

    content_slide(prs, "Does the tier mean something? Yes.", [
        "Checked against measures the tier was NOT built from:",
        "- SAT rises across every tier: 1,066 → 1,288, no inversions.",
        "- Graduation rate separates the bottom tier sharply (70% vs 84–89%).",
        "Honest note: the top tier reflects exam performance, not catalog size — so grad rate plateaus at the top rather than peaking.",
    ], image="rigor_validation.png",
       notes="The SAT check is the strongest single result: an independent measure, never used to "
             "build the tier, rises monotonically across all five. Volunteer the grad-rate dip "
             "before they find it.")

    content_slide(prs, "It predicts outcomes — and isn't just an SES proxy", [
        "Model: graduation rate ~ opportunity features (Adelman, 1999).",
        "Result: opportunity adds +5 pts R² beyond socioeconomic status — stable across models.",
        "Tension: the outcome IS SES-driven (free-lunch rate dominates)…",
        "…yet the index correlates with poverty at only −0.11 → opportunity, not laundered demographics.",
    ], image="predictive_validation.png",
       notes="This is the defensive slide. If anyone claims the tier just re-ranks wealthy "
             "suburbs, this is the answer -- and the poverty correlation is the number to quote.")

    content_slide(prs, "What we'd need from you, and how it ties in", [
        "Data ask: per-course AP score distributions (Calc BC vs Calc BC across schools) — lets us verify rigor at the course level.",
        "Open question: the measurement vintage of the org export — what year do those SAT/AP averages describe?",
        "Level: our tier is school-level; your counselor 'most demanding coursework' field is per-student — they complement each other.",
        "Handoff: the methodology is built to recalibrate against your internal outcomes (enrollment, GPA, retention).",
    ], notes="The real purpose of the meeting. Land the per-course ask concretely and get a yes/no. "
             "The vintage question has been open since Week 4.")

    content_slide(prs, "Next steps & discussion", [
        "Tier public and private separately — they aren't built from the same data.",
        "Validate on a second outcome (college-going).",
        "Package the methodology for handoff to your team.",
        "Your input: tier definitions, the per-course data ask, priority use-case (contextualization vs. outreach).",
    ], notes="End on their decision, not ours. The priority use-case question determines what we "
             "build in the remaining weeks.")
    out = os.path.join(ROOT, "docs", "Bob_Week6_Update.pptx")
    prs.save(out); return out


# ================================================================= MIDTERM DECK
def build_midterm():
    prs = new_deck()
    title_slide(prs,
        "An AI-Driven High School Rigor Classification",
        "MSDS 498 Capstone — Midterm  ·  Week 6",
        TEAM,
        notes="One-line pitch: admissions lost its standardized school-context tool, and we "
              "rebuilt one from public data that anyone can audit.")

    content_slide(prs, "Problem & objective", [
        "College Board discontinued Landscape (Sept 2025) — admissions lost standardized high-school context.",
        "~25% of applicants arrive with no school-context profile (Bastedo et al., 2023).",
        "The reader's problem: the same transcript means different things at different schools.",
        "Objective: a transparent, reproducible measure of a school's academic rigor — so a student's coursework can be read in the context of what their school actually offered.",
    ], notes="Be precise here: we measure schools, not students, and there is no admissions "
             "outcome in this data. The student-level version is the handoff, not this project.")

    content_slide(prs, "Data & reproducible pipeline", [
        "Sources: NCES, CRDC 2021-22, Census F-33/SAIPE, NU export, IB, ISBE.",
        "Pipeline: Dockerized ETL, re-runnable without engineering support (Boettiger, 2015).",
        "Scale: ~34,000 high schools.",
        "Governance: per-variable data dictionary + source/vintage tracking (client requirement).",
    ], image="pipeline.png",
       notes="The pipeline is a deliverable in its own right -- the client can re-run it when new "
             "CRDC or NCES vintages land.")

    content_slide(prs, "Record linkage without a shared identifier", [
        "Problem: CEEB and NCES have no common key.",
        "Decision rule: auto-accept / review / reject = Fellegi–Sunter three-way (1969).",
        "Similarity: token-based over edit distance for school names (Cohen et al., 2003).",
        "Ambiguous middle band: an LLM adjudicates the review-tier pairs a threshold cannot settle — each decision logged and auditable.",
        "Reporting: match rates reported as a matter of course.",
    ], image="match_rates.png",
       notes="Classic entity-resolution problem. The three-way decision rule matters because "
             "forcing a binary match/no-match would quietly inject false links into everything "
             "downstream.")

    content_slide(prs, "Feature engineering → the rigor index", [
        "Form: a weighted composite index — transparent and decomposable, not a black box.",
        "Principle: performance over availability (Geiser & Santelices, 2004).",
        "Structure: five components, per-row proportional weighting over available components.",
        "Latest revision: AP qualifying density (expected passing exams per student) + a verified IB signal.",
    ], image="index_schematic.png",
       notes="Qualifying density replaced mean AP score because a mean rewards gatekeeping -- a "
             "school that only lets its strongest students sit the exam posts a high average.")

    content_slide(prs, "Nominal vs. effective weights", [
        "Issue: assigned weights ≠ actual influence when features correlate (CADRE, 2024).",
        "Finding: AP performance leads the effective weight (~0.31); participation is absorbed.",
        "Sensitivity: drop the performance components and 30% of schools change tier — the literature-driven change was not cosmetic.",
        "Practice: both reported, across four weighting schemes.",
    ], image="weights.png",
       notes="The 30% number is the answer to 'did following the literature actually matter?'. "
             "Equal weights move only 7%, so the index is stable to reasonable reweighting but "
             "genuinely sensitive to dropping performance.")

    content_slide(prs, "From score to tiers — natural breaks", [
        "Method: Jenks natural breaks (Jenks, 1967; Fisher, 1958) — 1-D k-means, cutting at gaps in the distribution.",
        "Why not quantiles: Reardon cautions against splitting near-identical schools.",
        "Frame: norm-referenced (Glaser, 1963); criterion-referenced is the alternative (Cizek & Bunch, 2007).",
        "Result: 'Most Demanding' = 295 schools, not a forced top fifth.",
    ], image="tier_cutpoints.png",
       notes="Key caveat to state out loud: these cuts are relative to our scored population, not "
             "an absolute standard. Re-run on a different population and they move.")

    content_slide(prs, "Results — what each tier actually looks like", [
        "21,951 schools scored; 12,441 left unscored rather than imputed.",
        "AP exam score rises 2.06 → 3.56 across tiers — an index input, so this is internal consistency.",
        "SAT rises 1,066 → 1,288 — not an input, so this is external validation.",
        "Graduation rate separates the bottom tier, then plateaus — reported as found.",
    ], image="tier_profile.png",
       notes="The results slide. Walk the four panels left to right: how many, how they differ on "
             "an input, then on two things the model never saw.")

    content_slide(prs, "What comes out, per school", [
        "Deliverable: tier, score, percentile, and the components behind it — keyed on CEEB.",
        "Decomposable by construction: no school gets a tier without a traceable reason.",
        "Missingness surfaced: '4 of 5 components' is part of the output, not a footnote.",
        "Face validity: known-selective schools land where domain experts expect.",
    ], image="school_profile.png",
       notes="Anticipate 'show me a school I know'. Having one on the slide beats improvising.")

    content_slide(prs, "Validation against independent outcomes", [
        "SAT rises across tiers (1,066 → 1,288, a 222-point spread), no inversions — and never a model input.",
        "Graduation rate sharply separates the bottom tier.",
        "Honest finding: grad rate plateaus at the top — 'most demanding' means performance, not catalog size.",
        "Face-validity audit across known schools as a qualitative check.",
    ], image="rigor_validation.png",
       notes="Strongest empirical result in the deck. Monotonic across all five tiers on a measure "
             "the index never touched.")

    content_slide(prs, "Why opportunity, not test scores", [
        "The obvious alternative: just rank schools by mean SAT. We tested that.",
        "SAT correlates with county child poverty at −0.385; our tier, at −0.110.",
        "So an outcome measure is ~3.5× more socioeconomically confounded than an opportunity measure.",
        "This is empirical support for the design choice, not just a literature citation (Reardon / SEDA).",
    ], image="benchmarking_ses.png",
       notes="This slide converts a literature argument into our own evidence. If you only keep "
             "one methods-defense slide, keep this one.")

    content_slide(prs, "Predictive validation — the supervised model", [
        "Design: HistGradientBoosting + linear baseline; grad rate ~ opportunity, SES-controlled (Adelman; Reardon).",
        "Discipline: 80/20 held-out split, fixed seed, permutation importance (10 repeats), R²/RMSE on the test set.",
        "Result: opportunity adds +0.05 R² beyond SES — stable across two specs, both model families.",
        "Importance: free-lunch rate dominates (0.53) — the outcome is SES-driven…",
        "…yet the index correlates with poverty only −0.11 → evidence it is not an SES proxy.",
    ], image="predictive_validation.png",
       notes="Emphasise that this model validates the index; it is not the deliverable. Predicting "
             "graduation is how we test that the construct carries signal.")

    content_slide(prs, "Unsupervised segments", [
        "K-means + hierarchical over PCA components (90% variance retained); region, academic profile, funding.",
        "k=4 selected by gap statistic (Tibshirani et al., 2001), cross-checked against silhouette; the two algorithms compared by adjusted Rand index.",
        "Deliberately excludes rigor_score — otherwise 'do the clusters match the tier?' is circular.",
        "Finding: every cluster spans several tiers — segments describe a different cut of the data.",
        "Limitation: 5,801 complete-case schools — the clustering universe is much smaller than the tiering one.",
    ], image="clustering.png",
       notes="Expect 'why not just cluster to get the tiers?'. Because clustering gives unordered "
             "groups, and rigor is inherently ordinal -- plus no ground truth to name the order.")

    content_slide(prs, "A three-layer modeling system", [
        "Measurement model: the composite rigor index — the deliverable.",
        "Unsupervised ML: K-means + hierarchical clustering with PCA — school segments.",
        "Supervised ML: gradient boosting predicting outcomes — validation, not the product.",
        "Why not train the tier directly: no ground-truth rigor label exists anywhere — inventing one would make the model unfalsifiable.",
        "So: an auditable measurement instrument, validated by ML rather than produced by it.",
    ], notes="This is the 'where is the model?' answer. Three layers, each doing a job the others "
             "cannot. The absence of a ground-truth label is the reason for the architecture.")

    content_slide(prs, "Limitations (documented, not hidden)", [
        "Ecological inference: we measure the school, not the student.",
        "Levels, not growth: Reardon's caution — we report the SES-confounding check.",
        "Coverage: performance data ~35%, skewed to the recruiting universe.",
        "Gaps: private-school blind spot (CRDC public-only); per-course AP scores unavailable.",
        "Fragility: losing CRDC access would reshuffle ~40% of the CRDC-covered population.",
    ], notes="Volunteering limitations is what buys credibility on the rest. Each one maps to an "
             "extension point in the NU handoff rather than a dead end.")

    content_slide(prs, "Next steps (Weeks 7 → 10)", [
        "Separate public/private tiering — different feature availability, different models.",
        "Second-outcome validation (college-going), and a post-COVID graduation vintage.",
        "Methodology handoff guide for NU integration.",
        "Final report: methods, results, limitations.",
    ], notes="Close on the handoff. The measure is built to be recalibrated against NU's internal "
             "outcomes, which is the natural continuation of the work.")
    out = os.path.join(ROOT, "docs", "MSDS_498_Midterm.pptx")
    prs.save(out); return out


if __name__ == "__main__":
    for path in (build_bob(), build_midterm()):
        n = len(Presentation(path).slides._sldIdLst)
        print(f"wrote {path}  ({n} slides)")
